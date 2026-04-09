"""
Gemini service for text extraction from files.

This service handles interaction with Google's Gemini AI model
to extract text from various file formats.

Gemini multimodal `generate_content` accepts only a limited set of document MIME types
(e.g. PDF, plain text), not Office Open XML (.docx / .pptx). Those are read locally.
"""

import asyncio
import mimetypes
import os
from pathlib import Path

import google.generativeai as genai
from docx import Document as DocxDocument
from google.generativeai.types import HarmBlockThreshold, HarmCategory
from lxml import html as lxml_html
from pptx import Presentation

from src.core.gemini_config import gemini_config
from src.core.logging import get_logger

logger = get_logger(__name__)

# Formats we extract without calling Gemini (unsupported or unnecessary for the Files API).
_LOCAL_TEXT_EXTRACTION_EXTENSIONS: frozenset[str] = frozenset({
    ".docx",
    ".pptx",
    ".txt",
    ".html",
})

# google.generativeai.upload_file does not always infer MIME type from the path;
# map extensions we allow in config so uploads work with any filename (e.g. non-Latin).
_MIME_BY_EXTENSION: dict[str, str] = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
    ".heic": "image/heic",
    ".heif": "image/heif",
    ".pdf": "application/pdf",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".txt": "text/plain",
    ".html": "text/html",
}


def _resolve_upload_mime_type(file_name: str, content_type: str | None) -> str:
    if content_type:
        primary = content_type.split(";")[0].strip()
        if primary and primary.lower() != "application/octet-stream":
            return primary
    guessed, _ = mimetypes.guess_type(file_name)
    if guessed:
        return guessed
    ext = os.path.splitext(file_name)[1].lower()
    if ext in _MIME_BY_EXTENSION:
        return _MIME_BY_EXTENSION[ext]
    raise ValueError(
        "Could not determine MIME type for this file; "
        "provide a correct Content-Type or a file name with a known extension."
    )


def _extract_docx_text(path: str) -> str:
    doc = DocxDocument(path)
    parts: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for cp in cell.paragraphs:
                    t = cp.text.strip()
                    if t:
                        parts.append(t)
    return "\n".join(parts)


def _extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    parts.append(t)
    return "\n".join(parts)


def _extract_txt_or_html(path: str, ext: str) -> str:
    raw = Path(path).read_bytes()
    if ext == ".html":
        tree = lxml_html.fromstring(raw)
        return tree.text_content()
    return raw.decode("utf-8", errors="replace")


def _extract_text_locally(file_path: str, ext: str) -> str:
    if ext == ".docx":
        return _extract_docx_text(file_path)
    if ext == ".pptx":
        return _extract_pptx_text(file_path)
    if ext in (".txt", ".html"):
        return _extract_txt_or_html(file_path, ext)
    raise ValueError(f"Local extraction not implemented for {ext!r}")


class GeminiTextExtractor:
    """Service for extracting text from files using Gemini."""

    def __init__(self):
        """Initialize Gemini text extractor."""
        genai.configure(api_key=gemini_config.GEMINI_API_KEY)
        self.model = genai.GenerativeModel(gemini_config.GEMINI_MODEL)
        self.safety_settings = {
            HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
        }

    async def extract_text_from_file(
        self,
        file_path: str,
        file_name: str,
        content_type: str | None = None,
    ) -> str:
        """
        Extract text from a file using Gemini.

        Args:
            file_path: Path to the temporary file
            file_name: Original name of the file
            content_type: MIME type from the client upload (strongly recommended)

        Returns:
            Extracted text content

        Raises:
            Exception: If text extraction fails
        """
        ext = os.path.splitext(file_name)[1].lower()

        if ext == ".doc":
            raise ValueError(
                "Legacy Word .doc files are not supported. "
                "Please upload .docx or PDF instead."
            )

        if ext in _LOCAL_TEXT_EXTRACTION_EXTENSIONS:
            try:
                logger.info(
                    "extracting_text_locally",
                    file_name=file_name,
                    extension=ext,
                )
                text = await asyncio.to_thread(_extract_text_locally, file_path, ext)
                text = (text or "").strip()
                if not text:
                    raise ValueError("No readable text found in document")
                logger.info(
                    "text_extraction_successful",
                    file_name=file_name,
                    text_length=len(text),
                    source="local",
                )
                return text
            except Exception as e:
                logger.error(
                    "text_extraction_failed",
                    file_name=file_name,
                    error=str(e),
                    error_type=type(e).__name__,
                    exc_info=True,
                )
                raise RuntimeError(f"Failed to extract text from file: {str(e)}")

        uploaded_file = None
        try:
            mime_type = _resolve_upload_mime_type(file_name, content_type)
            logger.info(
                "uploading_file_to_gemini",
                file_name=file_name,
                mime_type=mime_type,
            )

            # Upload file to Gemini
            uploaded_file = genai.upload_file(file_path, mime_type=mime_type)

            # Wait for file processing
            logger.info("waiting_for_file_processing", file_name=file_name)
            await self._wait_for_file_processing(uploaded_file.name)

            # Extract text using Gemini
            extraction_prompt = """
            Extract all text content from this document. 
            
            Rules:
                - Return ONLY plain text.
                - Do NOT include tables, table formatting, columns, bullet points, or markdown.
                - Do NOT describe layout or structure.
                - Do NOT add explanations, comments, or headings.
                - Output must be a continuous text suitable for machine learning input.
                - If no readable text is present, return an empty string.
                - Preserve the original structure and organization of the text as much as possible.
                
            Return only the extracted text.
            """

            response = self.model.generate_content(
                contents=[extraction_prompt, uploaded_file],
                safety_settings=self.safety_settings,
            )

            text = self.extract_text_safe(response)
            if not text or not text.strip():
                raise ValueError("No readable text found in document")

            out = text.strip()
            logger.info(
                "text_extraction_successful",
                file_name=file_name,
                text_length=len(out),
                source="gemini",
            )

            return out

        except Exception as e:
            logger.error(
                "text_extraction_failed",
                file_name=file_name,
                error=str(e),
                error_type=type(e).__name__,
                exc_info=True
            )
            raise RuntimeError(f"Failed to extract text from file: {str(e)}")

        finally:
            # Clean up uploaded file from Gemini
            if uploaded_file:
                try:
                    genai.delete_file(uploaded_file.name)
                    logger.debug("gemini_file_deleted", file_name=file_name)
                except Exception as e:
                    logger.warning(
                        "failed_to_delete_gemini_file",
                        file_name=file_name,
                        error=str(e)
                    )

    async def _wait_for_file_processing(self, file_name: str, max_retries: int = 10, retry_delay: int = 5):
        """
        Wait for uploaded file to be processed by Gemini.

        Args:
            file_name: Name of the uploaded file in Gemini
            max_retries: Maximum number of retries
            retry_delay: Delay between retries in seconds

        Raises:
            Exception: If file processing times out
        """
        for attempt in range(max_retries):
            try:
                file_info = genai.get_file(file_name)
                # State 2 = PROCESSED, State 1 = PENDING
                if file_info.state == 2:
                    logger.debug(
                        "file_processing_complete",
                        file_name=file_name,
                        attempt=attempt + 1
                    )
                    return

                logger.info(
                    "file_processing_pending",
                    file_name=file_name,
                    state=file_info.state,
                    attempt=attempt + 1,
                    retry_in=retry_delay
                )
                await asyncio.sleep(retry_delay)

            except Exception as e:
                logger.warning(
                    "error_checking_file_status",
                    file_name=file_name,
                    attempt=attempt + 1,
                    error=str(e)
                )
                await asyncio.sleep(retry_delay)

        raise Exception(f"File processing timed out after {max_retries * retry_delay} seconds")

    @staticmethod
    def extract_text_safe(response) -> str | None:
        if not response.candidates:
            return None

        parts = response.candidates[0].content.parts
        texts = [
            part.text
            for part in parts
            if hasattr(part, "text") and part.text.strip()
        ]

        return "\n".join(texts) if texts else None